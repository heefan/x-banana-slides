"""
Export Service - handles PPTX and PDF export
Based on demo.py create_pptx_from_images()
"""
import os
import logging
import tempfile
from pathlib import Path
from typing import List, Optional, Dict
from contextlib import contextmanager
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
import io

logger = logging.getLogger(__name__)


class ExportService:
    """Service for exporting presentations"""
    
    @staticmethod
    def create_pptx_from_images(image_paths: List[str], output_file: str = None) -> bytes:
        """
        Create PPTX file from image paths
        Based on demo.py create_pptx_from_images()
        
        Args:
            image_paths: List of absolute paths to images
            output_file: Optional output file path (if None, returns bytes)
        
        Returns:
            PPTX file as bytes if output_file is None
        """
        # Create presentation
        prs = Presentation()
        
        # Set slide dimensions to 16:9 (width 10 inches, height 5.625 inches)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        
        # Add each image as a slide
        for image_path in image_paths:
            if not os.path.exists(image_path):
                logger.warning(f"Image not found: {image_path}")
                continue
            
            # Add blank slide layout (layout 6 is typically blank)
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add image to fill entire slide
            slide.shapes.add_picture(
                image_path,
                left=0,
                top=0,
                width=prs.slide_width,
                height=prs.slide_height
            )
        
        # Save or return bytes
        if output_file:
            prs.save(output_file)
            return None
        else:
            # Save to bytes
            pptx_bytes = io.BytesIO()
            prs.save(pptx_bytes)
            pptx_bytes.seek(0)
            return pptx_bytes.getvalue()
    
    @staticmethod
    def create_pdf_from_images(image_paths: List[str], output_file: str = None) -> bytes:
        """
        Create PDF file from image paths
        
        Args:
            image_paths: List of absolute paths to images
            output_file: Optional output file path (if None, returns bytes)
        
        Returns:
            PDF file as bytes if output_file is None
        """
        images = []
        
        # Load all images
        for image_path in image_paths:
            if not os.path.exists(image_path):
                logger.warning(f"Image not found: {image_path}")
                continue
            
            img = Image.open(image_path)
            
            # Convert to RGB if necessary (PDF requires RGB)
            if img.mode != 'RGB':
                img = img.convert('RGB')
            
            images.append(img)
        
        if not images:
            raise ValueError("No valid images found for PDF export")
        
        # Save as PDF
        if output_file:
            images[0].save(
                output_file,
                save_all=True,
                append_images=images[1:],
                format='PDF'
            )
            return None
        else:
            # Save to bytes
            pdf_bytes = io.BytesIO()
            images[0].save(
                pdf_bytes,
                save_all=True,
                append_images=images[1:],
                format='PDF'
            )
            pdf_bytes.seek(0)
            return pdf_bytes.getvalue()
    
    @staticmethod
    @contextmanager
    def _temp_image_file(image: Image.Image):
        """
        Temporary file context manager for image files
        
        Args:
            image: PIL Image object
            
        Yields:
            Temporary file path
        """
        temp_file = None
        try:
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as f:
                temp_file = f.name
                # Ensure image is in correct mode
                if image.mode == 'RGBA':
                    # If transparency is needed, use PNG
                    image.save(temp_file, format='PNG')
                else:
                    # Otherwise convert to RGB to save space
                    if image.mode != 'RGB':
                        image = image.convert('RGB')
                    image.save(temp_file, format='PNG')
            yield temp_file
        finally:
            if temp_file and os.path.exists(temp_file):
                try:
                    os.remove(temp_file)
                except Exception as e:
                    logger.warning(f"Failed to remove temp file {temp_file}: {e}")
    
    @staticmethod
    def _add_slide_with_elements(
        slide,
        image_path: str,
        segmentation_service,
        prs: Presentation,
        original_image: Image.Image
    ):
        """
        Add a slide with segmented elements
        
        Args:
            slide: PPTX slide object
            image_path: Path to original image
            segmentation_service: ElementSegmentationService instance
            prs: Presentation object
            original_image: PIL Image object of the original image
        """
        # Get image dimensions
        image_width, image_height = original_image.size
        
        # Calculate scale factors (EMU per pixel)
        # python-pptx uses EMU internally (1 inch = 914400 EMU)
        scale_x = prs.slide_width / image_width  # EMU per pixel
        scale_y = prs.slide_height / image_height  # EMU per pixel
        
        # Step 1: Segment image to get elements (before adding background)
        try:
            elements_info = segmentation_service.segment_image(image_path)
        except Exception as e:
            logger.warning(f"Segmentation failed for {image_path}: {e}, will use simple export")
            # Raise exception so caller can handle fallback
            raise
        
        # Step 2: Add background image (bottom layer)
        slide.shapes.add_picture(
            image_path,
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height
        )
        
        # Step 3: Add text elements (middle layer)
        for text_elem in elements_info.get('text_elements', []):
            try:
                bbox = text_elem.get('bbox', [])
                if len(bbox) != 4:
                    logger.warning(f"Invalid bbox for text element: {bbox}")
                    continue
                
                text_content = text_elem.get('text', '').strip()
                if not text_content:
                    continue
                
                # Convert pixel coordinates to PPT EMU coordinates
                left_emu = bbox[0] * scale_x
                top_emu = bbox[1] * scale_y
                width_emu = bbox[2] * scale_x
                height_emu = bbox[3] * scale_y
                
                # Create textbox (convert EMU to Inches)
                textbox = slide.shapes.add_textbox(
                    Inches(left_emu / 914400),
                    Inches(top_emu / 914400),
                    Inches(width_emu / 914400),
                    Inches(height_emu / 914400)
                )
                
                text_frame = textbox.text_frame
                text_frame.text = text_content
                text_frame.word_wrap = True
                
                # Set font style (precise conversion)
                font_size_px = text_elem.get('font_size', 18)
                
                # Calculate font size in points
                # Formula: font_size_pt = font_size_px * (ppt_height_inches / image_height_px) * 72
                ppt_height_inches = prs.slide_height / 914400
                font_size_pt = font_size_px * (ppt_height_inches / image_height) * 72
                
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = Pt(int(font_size_pt))
                    paragraph.font.bold = (text_elem.get('font_weight') == 'bold')
                
                logger.debug(f"Added text element: {text_content[:50]}... at {bbox}")
                
            except Exception as e:
                logger.warning(f"Failed to add text element: {e}", exc_info=True)
                continue
        
        # Step 4: Add icons and charts (top layer)
        all_visual_elements = elements_info.get('icons', []) + elements_info.get('charts', [])
        
        for elem in all_visual_elements:
            try:
                bbox = elem.get('bbox', [])
                if len(bbox) != 4:
                    logger.warning(f"Invalid bbox for visual element: {bbox}")
                    continue
                
                # Crop element from original image
                icon_image = original_image.crop((
                    int(bbox[0]),
                    int(bbox[1]),
                    int(bbox[0] + bbox[2]),
                    int(bbox[1] + bbox[3])
                ))
                
                # Use temporary file context manager
                with ExportService._temp_image_file(icon_image) as temp_path:
                    # Convert pixel coordinates to PPT EMU coordinates
                    left_emu = bbox[0] * scale_x
                    top_emu = bbox[1] * scale_y
                    width_emu = bbox[2] * scale_x
                    height_emu = bbox[3] * scale_y
                    
                    # Add picture (convert EMU to Inches)
                    slide.shapes.add_picture(
                        temp_path,
                        left=Inches(left_emu / 914400),
                        top=Inches(top_emu / 914400),
                        width=Inches(width_emu / 914400),
                        height=Inches(height_emu / 914400)
                    )
                
                logger.debug(f"Added visual element at {bbox}")
                
            except Exception as e:
                logger.warning(f"Failed to add visual element: {e}", exc_info=True)
                continue
    
    @staticmethod
    def create_pptx_with_segmented_elements(
        image_paths: List[str],
        output_file: str = None,
        use_segmentation: bool = True
    ) -> Optional[bytes]:
        """
        Create PPTX file with segmented elements (editable text and movable icons)
        
        Args:
            image_paths: List of absolute paths to images
            output_file: Optional output file path (if None, returns bytes)
            use_segmentation: Whether to use element segmentation (default: True)
        
        Returns:
            PPTX file as bytes if output_file is None, None otherwise
        
        Note:
            If use_segmentation=True but segmentation fails, automatically falls back
            to simple export (whole image)
        """
        # Initialize segmentation service
        segmentation_service = None
        if use_segmentation:
            try:
                from services.element_segmentation_service import ElementSegmentationService
                segmentation_service = ElementSegmentationService()
                logger.info("ElementSegmentationService initialized successfully")
            except Exception as e:
                logger.warning(f"Failed to initialize segmentation: {e}, using simple export")
                use_segmentation = False
        
        # Create presentation
        prs = Presentation()
        
        # Set slide dimensions to 16:9 (width 10 inches, height 5.625 inches)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        
        # Process each image
        for image_path in image_paths:
            if not os.path.exists(image_path):
                logger.warning(f"Image not found: {image_path}")
                continue
            
            # Add blank slide layout
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Load original image for cropping
            try:
                original_image = Image.open(image_path)
            except Exception as e:
                logger.error(f"Failed to load image {image_path}: {e}")
                # Fallback to simple export
                slide.shapes.add_picture(
                    image_path,
                    left=0,
                    top=0,
                    width=prs.slide_width,
                    height=prs.slide_height
                )
                continue
            
            if use_segmentation and segmentation_service:
                try:
                    # Try to use element segmentation
                    ExportService._add_slide_with_elements(
                        slide,
                        image_path,
                        segmentation_service,
                        prs,
                        original_image
                    )
                    logger.info(f"Successfully added slide with segmented elements: {image_path}")
                except Exception as e:
                    logger.warning(
                        f"Segmentation failed for {image_path}: {e}, falling back to simple export",
                        exc_info=True
                    )
                    # Fallback to simple export
                    slide.shapes.add_picture(
                        image_path,
                        left=0,
                        top=0,
                        width=prs.slide_width,
                        height=prs.slide_height
                    )
            else:
                # Direct simple export
                slide.shapes.add_picture(
                    image_path,
                    left=0,
                    top=0,
                    width=prs.slide_width,
                    height=prs.slide_height
                )
        
        # Save or return bytes
        if output_file:
            prs.save(output_file)
            return None
        else:
            # Save to bytes
            pptx_bytes = io.BytesIO()
            prs.save(pptx_bytes)
            pptx_bytes.seek(0)
            return pptx_bytes.getvalue()

